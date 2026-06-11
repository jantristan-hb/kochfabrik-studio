"""US-012 — Konformitäts-Check für aus Template+Modell erzeugte Angebote.

Akzeptanzkriterium 1 (Vorstufe zum Pixel-Diff-Gate in Sprint 3): ein
aus angebot_template + angebot_model befülltes, gerendertes Deck muss
strukturell ein echtes KOCHfabrik-Angebot sein.

Ablauf: cached_deck (Assets/Logos materialisieren) → angebot_fill.fill
→ reconstruct.js → Text aus dem gerenderten PPTX → kf_classify.
(PDF→pdftotext = Sprint-3-Pipeline; kf_classify ist textbasiert, der
Check auf dem gerenderten PPTX-Text ist äquivalent — Sprint-2-Scope.)

Run: python3 verify_angebot.py
"""
import json
import os
import re
import subprocess
import sys
import tempfile

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
# US-056: Tooling lebt jetzt unter engine/tooling/ — Runtime-Module
# (engine/scripts/) zusätzlich auf den Pfad (KEINE Logikänderung).
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                "..", "scripts"))
from _deckpipe import cached_deck                               # noqa
from angebot_fill import fill, TEMPLATE                          # noqa
from angebot_model import example                                # noqa
from kf_classify import is_kochfabrik, classify, extract_event   # noqa

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SPIKE = os.path.join(ROOT, "spike-pptxgenjs")
SCRIPTS = os.path.join(ROOT, "scripts")
REF_PDF = ("/Users/janrudat/Nextcloud/Kochfabrik Dokumente/"
           "AKARA_Muster_Angebote/# 10_182_RAUMKARUSSELL GmbH_"
           "12_09_2026.pdf")
LABELS = ["Veranstaltungsanlass", "Veranstaltungsdatum",
          "Veranstaltungsbeginn", "Personenanzahl",
          "Veranstaltungsort", "Cateringkonzept"]


def _pptx_text(path):
    from pptx import Presentation
    out = []
    for sl in Presentation(path).slides:
        for sh in sl.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                out.append(sh.text_frame.text)
    return re.sub(r"\s+", " ", " ".join(out))


def verify():
    if not os.path.isfile(TEMPLATE):
        subprocess.run([sys.executable,
                        os.path.join(SCRIPTS, "build_angebot_template.py")],
                       check=True, capture_output=True)
    shared = tempfile.mkdtemp(prefix="angverify_")
    _, _, logos = cached_deck(REF_PDF, shared)            # Assets in shared
    a = example()
    el = fill(a)
    json.dump(el, open(os.path.join(shared, "elements.json"), "w"),
              ensure_ascii=False)
    json.dump(logos, open(os.path.join(shared, "logos.json"), "w"))
    pptx = os.path.join(shared, "angebot_verify.pptx")
    r = subprocess.run(["node", os.path.join(SPIKE, "reconstruct.js"),
                        "elements.json", pptx], cwd=shared,
                       capture_output=True, text=True)
    txt = _pptx_text(pptx) if r.returncode == 0 else ""
    ev = extract_event(txt)

    checks = {
        "reconstruct rc==0": r.returncode == 0,
        "is_kochfabrik": is_kochfabrik(txt),
        "classify == angebot": classify(txt, 0) == "angebot",
        "alle 6 Labels": all(L in txt for L in LABELS),
        "Bankblock (BIC/Standorte)":
            "GENODEF1PIN" in txt or "Planungsfabrik" in txt,
        "Kundenwert injiziert (RAUMKARUSSELL)": a.kunde in txt,
        "Anlass injiziert": a.veranstaltung.anlass in txt,
    }
    return all(checks.values()), checks, r.stderr[-200:]


if __name__ == "__main__":
    ok, checks, err = verify()
    for k, v in checks.items():
        print(("  OK  " if v else "  FAIL") + " " + k)
    if not ok and err:
        print("stderr:", err)
    print("KONFORM" if ok else "NICHT KONFORM")
    sys.exit(0 if ok else 1)
