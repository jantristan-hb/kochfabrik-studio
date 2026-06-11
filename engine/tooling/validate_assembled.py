"""validate_assembled.py — LOKALER WEGWERF-Validator: assembled-Deck
↔ Angebot. Kein DB/API. Zum Algo-Feintuning, nicht für Produktion.

Prüft: Slide-Zahl plausibel? #Food == #Gänge? Cover/Frame/Ausstattung
vollständig + Reihenfolge (Cover zuerst, Kontakt zuletzt)? Pro Gang:
Gericht-Abdeckung (Token-Overlap Angebot↔Slide) + Overflow ('•'-Run-on).

Usage:
  python3 validate_assembled.py <offer.(md|pdf)> [--offer SEL] <deck.pptx>
"""
import argparse
import os
import re
import sys

from pptx import Presentation

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
# US-056: Tooling lebt jetzt unter engine/tooling/ — Runtime-Module
# (engine/scripts/) zusätzlich auf den Pfad (KEINE Logikänderung).
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                "..", "scripts"))
from compose_offer import parse_offer_dishes                   # noqa

FRAME = {"KONTAKT", "WERTSCHÄTZUNG IST DER SCHLÜSSEL",
         "DEINE CATERING- & EVENT-CREW IM NORDEN", "PERSONAL",
         "AUSTATTUNG", "AUSSTATTUNG", "LOCATION_AUSSTATTUNG"}
STOP = set("UND ODER MIT IM IN DER DIE DAS VOM AUF & / | - "
           "AUS ZUM ZUR EIN EINE A AB".split())


def toks(s):
    return {w for w in re.findall(r"[a-zäöüß]+", (s or "").lower())
            if len(w) > 3 and w.upper() not in STOP}


def slide_txt(sl):
    big, bs, full = "", -1, []
    for sh in sl.shapes:
        if not sh.has_text_frame or not sh.text_frame.text.strip():
            continue
        t = sh.text_frame.text
        full.append(t)
        r = sh.text_frame.paragraphs[0].runs
        sz = (r[0].font.size.pt if r and r[0].font.size else 0) or 0
        if sz > bs:
            bs, big = sz, t.replace("\n", " ")
    return re.sub(r"\s+", " ", big).strip().upper(), " ".join(full)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("offer")
    ap.add_argument("deck")
    ap.add_argument("--offer", dest="sel", default="")
    a = ap.parse_args()

    courses = parse_offer_dishes(a.offer, a.sel)          # [(name,[(n,d)])]
    nC = len(courses)
    sl = list(Presentation(a.deck).slides)
    heads = [slide_txt(s) for s in sl]
    cover = heads[0][0] if heads else ""
    frame_h = [h for h, _ in heads if any(h.startswith(f[:14])
               for f in FRAME)]
    food = [(h, f) for h, f in heads[1:]
            if not any(h.startswith(x[:14]) for x in FRAME)]
    food = food[:nC] if len(food) >= nC else food          # erwartete Food

    exp = 1 + nC + len(frame_h) + 0       # +Ausstattung ist i.d.R. in frame
    w = 0

    def chk(c, name, extra=""):
        nonlocal w
        w += 0 if c else 1
        print(f"  [{'OK ' if c else 'WARN'}] {name}{extra}")

    print(f"== {os.path.basename(a.deck)}  (Angebot: {a.sel or a.offer}) ==")
    print(f"  Gänge={nC}  Slides={len(sl)}  Frame={len(frame_h)}  "
          f"Food={len(food)}")
    chk(len(sl) >= 1 + nC + 3, "Slide-Zahl plausibel",
        f"  (>= 1+{nC}+Frame; ist {len(sl)})")
    chk(len(food) == nC, "#Food == #Gänge", f"  ({len(food)}/{nC})")
    chk(bool(cover) and len(cover) > 6, "Cover mit Titel", f"  '{cover[:30]}'")
    for need in ("KONTAKT", "WERTSCHÄTZUNG", "DEINE CATERING", "PERSONAL"):
        chk(any(h.startswith(need[:14]) for h in frame_h),
            f"Frame vorhanden: {need[:18]}")
    chk(heads and heads[-1][0].startswith("KONTAKT"),
        "Kontakt = letzte Slide")
    # Gericht-Abdeckung je Gang (lexikalischer Proxy)
    lows = []
    for (cname, dishes), (fh, ff) in zip(courses, food):
        otok = set().union(*[toks(n + " " + d) for n, d in dishes]) \
            if dishes else set()
        stok = toks(ff)
        ov = len(otok & stok) / max(len(otok), 1)
        run = ff.count(" • ")
        flag = "OVERFLOW" if run >= 1 else ("LOW" if ov < 0.15 else "")
        if flag:
            lows.append((cname, fh, round(ov, 2), flag))
        print(f"   Gang «{cname[:20]:20}» → «{fh[:22]:22}» "
              f"overlap={ov:.0%} run-on={run} {flag}")
    chk(not lows, "Gericht-Abdeckung ok (kein LOW/OVERFLOW)",
        f"  ({len(lows)} schwach)")
    print(f"  → {'PASS' if w == 0 else str(w)+' WARN'}\n")
    sys.exit(1 if w else 0)


if __name__ == "__main__":
    main()
