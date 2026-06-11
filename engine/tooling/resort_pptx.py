"""resort_pptx.py — pptx nach Cluster-Tags umsortieren (in-place, gehärtet).

Liest <pptx>.tags.json (no,deck,page,headline,cluster) und ordnet die
Slides so um, dass semantisch gleiche (gleicher Cluster) am Stück liegen:
  - Cluster nach Größe absteigend (große Module zuerst)
  - innerhalb Cluster nach deck, page (gleiches Deck beisammen)
  - Singletons ans Ende, alphabetisch nach Headline (ähnliche nah)
Notizen 'deck::page' reisen mit dem Slide-Part → Ground-Truth bleibt.
Save härtet wie curate.py: dangling slide-Rels (Alt-Waisen) werden
mitkompaktiert.

Usage: python3 resort_pptx.py /tmp/all_menus.pptx
"""
import json
import os
import sys
from collections import Counter

from pptx import Presentation
from pptx.oxml.ns import qn

SLIDE_RT = ("http://schemas.openxmlformats.org/officeDocument/2006/"
            "relationships/slide")


def main():
    pptx = sys.argv[1] if len(sys.argv) > 1 else "/tmp/all_menus.pptx"
    tags = json.load(open(pptx.replace(".pptx", ".tags.json")
                          if pptx.endswith(".pptx")
                          else pptx + ".tags.json"))
    by_no = {r["no"]: r for r in tags}
    sizes = Counter(r["cluster"] for r in tags)

    def key(r):
        c = r["cluster"]
        s = sizes[c]
        # Multi-Cluster: nach Größe desc, dann Cluster-ID stabil, dann
        # deck/page. Singletons (s==1): Block ans Ende, Headline-alpha.
        if s > 1:
            return (0, -s, c, r["deck"], _pg(r["page"]))
        return (1, 0, 0, r["headline"], r["deck"], _pg(r["page"]))

    def _pg(p):
        try:
            return int(p)
        except (TypeError, ValueError):
            return 0

    prs = Presentation(pptx)
    lst = prs.slides._sldIdLst
    sld_ids = list(lst)                       # aktuelle Reihenfolge (1-based)
    n = len(sld_ids)
    order = sorted(range(1, n + 1),
                   key=lambda i: key(by_no[i]) if i in by_no else (2,))
    # sldId-Elemente in neuer Reihenfolge neu anhängen
    for e in sld_ids:
        lst.remove(e)
    for i in order:
        lst.append(sld_ids[i - 1])

    # Kompaktierung: slide-Rels ohne sldId-Eintrag (Alt-Waisen) wegwerfen
    live = {s.get(qn("r:id")) for s in lst}
    dropped = 0
    for rel in list(prs.part.rels.values()):
        if rel.reltype == SLIDE_RT and rel.rId not in live:
            prs.part.drop_rel(rel.rId)
            dropped += 1

    tmp = pptx + ".tmp"
    prs.save(tmp)
    os.replace(tmp, pptx)

    chk = Presentation(pptx)
    cs = list(chk.slides)
    wn = sum(1 for s in cs if s.has_notes_slide
             and s.notes_slide.notes_text_frame.text.strip())
    first = [by_no[i]["headline"][:30] for i in order[:6] if i in by_no]
    print(f"Umsortiert: {n} Slides | Alt-Waisen kompaktiert: {dropped}")
    print(f"Verify: {len(cs)} Slides, {wn} mit deck::page-Notiz")
    print(f"Reihenfolge-Start: {first}")
    print(f"Dateigröße: {os.path.getsize(pptx) // (1024 * 1024)} MB")


if __name__ == "__main__":
    main()
