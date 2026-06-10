"""slide_text.py — pro Slide (deck,page,headline,body) → JSON.

Liest die Notiz 'deck::page' (Ground-Truth-Key) + Headline (größter Text)
+ zusammengefassten Body-Text. Read-only auf der pptx. Output ist der
Input fürs Embedding/Clustering.

Usage: python3 ../scripts/slide_text.py /tmp/all_menus.pptx [out.json]
"""
import json
import re
import sys

from pptx import Presentation

sys.path.insert(0, __file__.rsplit("/", 1)[0])
from curate import slide_texts  # noqa: E402


def main():
    inp = sys.argv[1] if len(sys.argv) > 1 else "/tmp/all_menus.pptx"
    out = sys.argv[2] if len(sys.argv) > 2 else inp + ".slides.json"
    prs = Presentation(inp)
    rows, no_note, no_head = [], 0, 0
    for i, sl in enumerate(prs.slides, 1):
        note = (sl.notes_slide.notes_text_frame.text.strip()
                if sl.has_notes_slide else "")
        if note and "::" in note:
            deck, _, pg = note.rpartition("::")
        else:
            deck, pg, no_note = "", "", no_note + 1
        head, full = slide_texts(sl)
        if not head:
            no_head += 1
        # Body = voller Text ohne die Headline, auf 600 Zeichen begrenzt
        body = re.sub(re.escape(head), "", full, count=1).strip()[:600]
        rows.append({"no": i, "deck": deck, "page": pg,
                     "headline": head, "body": body})
    json.dump(rows, open(out, "w"), ensure_ascii=False, indent=0)
    print(f"{len(rows)} Slides → {out}")
    print(f"  ohne Notiz: {no_note} | ohne Headline: {no_head}")
    uniq = len({r["headline"] for r in rows})
    print(f"  distinkte Headlines: {uniq}")


if __name__ == "__main__":
    main()
