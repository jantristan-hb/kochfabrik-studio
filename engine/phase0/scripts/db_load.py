"""db_load.py — kuratierte Ground-Truth → menu_composition.

Quelle der Wahrheit: überlebende 'deck::page'-Notizen in
phase0/data/all_menus.pptx (1010). Join:
  slides.json → headline, body
  tags.json   → cluster (= module_type); module_label = häufigste
                Headline im Cluster über die überlebenden Slides
  CORPUS_DIR  → src_pdf via Reverse-Slug (_deckpipe.slugify)
Embedding bleibt NULL (db_embed.py füllt vector(768), Step 4).

Usage: python3 db_load.py            # idempotent (TRUNCATE + reload)
"""
import json
import os
import sys
from collections import Counter

import psycopg2
from psycopg2.extras import execute_values

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from _deckpipe import slugify                                  # noqa

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DATA = os.path.join(ROOT, "data")
CORPUS = "/Users/janrudat/Nextcloud/Kochfabrik Dokumente/AKARA_Präsentationen"
DSN = dict(host="localhost", port=5434, user="postgres",
           password="pptxgen", dbname="pptxgen")


def main():
    from pptx import Presentation
    keep = set()
    for s in Presentation(os.path.join(DATA, "all_menus.pptx")).slides:
        if s.has_notes_slide:
            n = s.notes_slide.notes_text_frame.text.strip()
            if "::" in n:
                keep.add(n)
    print(f"Ground-Truth: {len(keep)} überlebende deck::page")

    slides = {f"{r['deck']}::{r['page']}": r
              for r in json.load(open(os.path.join(DATA,
                                                   "all_menus.slides.json")))}
    tags = {f"{r['deck']}::{r['page']}": r
            for r in json.load(open(os.path.join(DATA,
                                                 "all_menus.tags.json")))}
    # module_label = häufigste Headline je Cluster (nur überlebende)
    cl_heads = {}
    for k in keep:
        t = tags.get(k)
        if t:
            cl_heads.setdefault(t["cluster"], []).append(
                t.get("headline", ""))
    label = {c: Counter(h).most_common(1)[0][0] for c, h in cl_heads.items()}

    smap = {slugify(p): os.path.join(CORPUS, p)
            for p in os.listdir(CORPUS) if p.lower().endswith(".pdf")}

    rows, miss_src, miss_meta = [], 0, 0
    for k in sorted(keep):
        deck, _, pg = k.rpartition("::")
        sl, tg = slides.get(k), tags.get(k)
        if not sl:
            miss_meta += 1
            continue
        src = smap.get(deck)
        if not src:
            miss_src += 1
            continue
        cl = tg["cluster"] if tg else None
        rows.append((deck, src, int(pg), sl.get("headline", ""),
                     sl.get("body", ""), cl,
                     label.get(cl) if cl is not None else None))

    cx = psycopg2.connect(**DSN)
    cu = cx.cursor()
    cu.execute("TRUNCATE menu_composition RESTART IDENTITY")
    execute_values(cu,
                   "INSERT INTO menu_composition "
                   "(deck,src_pdf,page,headline,body,module_type,"
                   "module_label) VALUES %s", rows)
    cx.commit()
    cu.execute("SELECT count(*), count(distinct deck), "
               "count(distinct module_type) FROM menu_composition")
    n, d, m = cu.fetchone()
    print(f"Geladen: {n} Zeilen | {d} Decks | {m} Module")
    print(f"  ohne src_pdf übersprungen: {miss_src} | "
          f"ohne slides.json-Meta: {miss_meta}")
    cu.execute("SELECT module_label, count(*) c FROM menu_composition "
               "GROUP BY module_label ORDER BY c DESC LIMIT 8")
    print("  Top-Module:")
    for lab, c in cu.fetchall():
        print(f"    {c:4d}  {lab}")
    cx.close()


if __name__ == "__main__":
    main()
