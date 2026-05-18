"""readback_overrides.py — Hand-Korrekturen aus einem in PowerPoint/
LibreOffice editierten PPTX zurücklesen und pro DECK in overrides.json
persistieren.

Erfasst NUR Panel-artige Rechtecke (sizable, kein Voll-BG, kein dünner
Frame-Streifen) — genau die Elemente, die man von Hand verschiebt
(Banner/Bänder). Frame/Voll-BG/Hairlines bleiben unangetastet, damit die
Bleed-/z-Order-Logik der Engine nicht doppelt greift.

Usage: readback_overrides.py <edited.pptx> <deck-key> [overrides.json]
"""
import json
import os
import sys
from pptx import Presentation
from pptx.util import Emu


def inch(v):
    return round(Emu(v).inches, 3)


def main():
    if len(sys.argv) < 3:
        sys.exit("Usage: readback_overrides.py <edited.pptx> <deck-key> "
                 "[overrides.json]")
    pptx, deck = sys.argv[1], sys.argv[2]
    ovp = sys.argv[3] if len(sys.argv) > 3 else "overrides.json"

    pres = Presentation(pptx)
    SW, SH = Emu(pres.slide_width).inches, Emu(pres.slide_height).inches

    deck_map = {}
    for i, slide in enumerate(pres.slides, start=1):
        ents = []
        for sh in slide.shapes:
            try:
                rgb = str(sh.fill.fore_color.rgb)
            except Exception:
                continue
            x, y = inch(sh.left), inch(sh.top)
            w, h = inch(sh.width), inch(sh.height)
            full_bg = w >= 0.97 * SW and h >= 0.97 * SH
            thin = w <= 0.3 or h <= 0.3
            sizable = w >= 1.0 or h >= 0.5
            if full_bg or thin or not sizable:
                continue
            ents.append({
                "match": {"t": "rect", "fill": rgb,
                          "nearX": x, "nearY": y},
                "set": {"x": x, "y": y, "w": w, "h": h},
            })
        if ents:
            deck_map[str(i)] = ents

    data = {}
    if os.path.isfile(ovp):
        try:
            data = json.load(open(ovp))
        except Exception:
            data = {}
    # Legacy-Flat-Format (Seiten-Keys auf Top-Level) in Default-Deck heben
    if data and all(k.isdigit() for k in data):
        data = {"_legacy": data}
    data[deck] = deck_map
    json.dump(data, open(ovp, "w"), indent=2)
    print(f"overrides.json: deck '{deck}' aktualisiert "
          f"({sum(len(v) for v in deck_map.values())} Panels, "
          f"{len(deck_map)} Seiten)")


if __name__ == "__main__":
    main()
