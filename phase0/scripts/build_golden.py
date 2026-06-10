"""build_golden.py — MASTER-GOLDEN-DATASET der kanonischen Static-Frames.

Frequenz NUR über normalisierten Text (Geometrie fragmentiert ein
Template über Bbox-Extraktions-Jitter → falsch). Es zählt das EXAKTE
Slide-Vorkommen über die pristine 1053-Nicht-Food.

Golden = Text-Typen mit count >= --min (Default 20 ≈ 10% der 199 Decks)
→ isoliert die scharfe Kante (4 dominante Typen 144–155, Rest ≤11).
Eine kanonische Instanz je Typ, sortiert in **Skelett-Reihenfolge**
(Crew früh → Personal → Wertschätzung → Kontakt zuletzt; sonst nach
Ø-Position via Headline-Keyword, Fallback Häufigkeit).

Output: phase0/data/master_golden.pptx (+ .manifest.json), persistent.
Usage: python3 build_golden.py [--min 20]
"""
import argparse
import glob
import json
import os
import re
import subprocess
import sys
import tempfile
from collections import defaultdict

from pptx import Presentation

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from _deckpipe import cached_deck, slugify                     # noqa
from dedup_exact import CACHE                                  # noqa

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DATA = os.path.join(ROOT, "data")
SPIKE = os.path.join(ROOT, "spike-pptxgenjs")
CORPUS = "/Users/janrudat/Nextcloud/Kochfabrik Dokumente/AKARA_Präsentationen"

# Skelett-Position (aus REPORT-structure.md Ø-Position); kleiner = früher
SKELETON = [("CREW IM NORDEN", 0.10), ("PARTNER IM NORDEN", 0.10),
            ("SO EMPFANGEN", 0.19), ("PERSONAL", 0.76),
            ("AUSTATTUNG", 0.78), ("AUSSTATTUNG", 0.78),
            ("WERTSCHÄTZUNG", 0.88), ("KONTAKT", 1.00)]


def text_sig(seq):
    t = " ".join(l.get("txt", "") for e in seq if e.get("t") == "text"
                 for l in e.get("lines", []))
    return re.sub(r"\s+", " ", t).strip().upper()


def skel_pos(nt):
    for kw, p in SKELETON:
        if kw in nt:
            return p
    return 0.5                                   # unbekannt → Mitte


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--min", type=int, default=20,
                    help="Mindest-Vorkommen (Default 20 ≈ 10% Decks)")
    a = ap.parse_args()

    pristine = sorted(glob.glob(os.path.join(DATA, "all_info_full_*.pptx")))
    if not pristine:
        sys.exit("Kein pristine all_info_full_*.pptx")
    src_pptx = pristine[-1]
    print(f"Pristine: {os.path.basename(src_pptx)}")

    grp = defaultdict(list)                       # norm_text -> [(slug,page)]
    for s in Presentation(src_pptx).slides:
        n = (s.notes_slide.notes_text_frame.text.strip()
             if s.has_notes_slide else "")
        if "::" not in n:
            continue
        slug, _, pg = n.rpartition("::")
        cp = os.path.join(CACHE, slug, "elements.json")
        if not os.path.isfile(cp):
            continue
        seq = json.load(open(cp)).get(str(int(pg)))
        if not seq:
            continue
        nt = text_sig(seq)
        if len(nt) >= 12:
            grp[nt].append((slug, int(pg)))

    reps = [(len(v), nt, v[0]) for nt, v in grp.items() if len(v) >= a.min]
    # Skelett-Reihenfolge: nach Ø-Position, bei Gleichstand häufigste zuerst
    reps.sort(key=lambda r: (skel_pos(r[1]), -r[0]))
    print(f"Golden-Typen (count >= {a.min}): {len(reps)}")
    for c, nt, _ in reps:
        print(f"  {c:4d}  pos={skel_pos(nt):.2f}  {nt[:54]}")
    if not reps:
        sys.exit("Nichts über Schwelle.")

    pdfs = {slugify(p): os.path.join(CORPUS, p)
            for p in os.listdir(CORPUS) if p.lower().endswith(".pdf")}
    shared = tempfile.mkdtemp(prefix="golden_")
    cache_el, logos, meta = {}, {}, None
    for c, nt, (slug, pg) in reps:
        if slug in cache_el:
            continue
        src = pdfs.get(slug)
        if not src:
            continue
        _, el, lg = cached_deck(src, shared)
        logos.update(lg)
        if meta is None:
            meta = el.get("_meta", {"w_pt": 960, "h_pt": 540})
        cache_el[slug] = el

    combined, notes, manifest, n = {}, {}, {}, 0
    for c, nt, (slug, pg) in reps:
        el = cache_el.get(slug)
        seq = el.get(str(pg)) if el else None
        if not seq:
            print(f"  warn fehlt {slug}::{pg}", file=sys.stderr)
            continue
        n += 1
        combined[str(n)] = seq
        notes[str(n)] = f"{slug}::{pg}"
        manifest[str(n)] = {"count": c, "skel_pos": skel_pos(nt),
                            "type": nt[:90], "deck": slug, "page": pg}

    mm = dict(meta or {"w_pt": 960, "h_pt": 540})
    mm["deck"], mm["notes"] = "master-golden", notes
    combined["_meta"] = mm
    json.dump(combined, open(os.path.join(shared, "elements.json"), "w"))
    json.dump(logos, open(os.path.join(shared, "logos.json"), "w"))
    out = os.path.join(DATA, "master_golden.pptx")
    json.dump(manifest, open(out + ".manifest.json", "w"),
              ensure_ascii=False, indent=1)
    subprocess.run(["node", os.path.join(SPIKE, "reconstruct.js"),
                    "elements.json", out], cwd=shared,
                   capture_output=True, check=True, timeout=600)
    print(f"\nOK → {out} — {n} Golden-Slides in Skelett-Reihenfolge")


if __name__ == "__main__":
    main()
