"""db_load_static.py — static_slide aus der HAND-KURATION grounden.

Quelle der Wahrheit = Jans kuratierte Auswahl-pptx (default
phase0/data/category_samples.pptx) — NICHT Korpus-Frequenz (die
brachte Artefakte wie 'Crazy Kitchen'/'SEITE x/y' rein, die Jan nie
gewählt hat). In die Tabelle kommt NUR was in der Auswahl überlebt
hat.

Pro Slide: category = Headline, full_text, img_sig, (deck,page,src).
Gruppiert je Kategorie: erste = golden (is_golden), übrige =
is_golden=false (Jans freigegebene Foto-Set-Alternativen, Tier B).
cnt + skel_pos = Korpus-Metadaten (Häufigkeit/Ø-Position des
exakten Volltexts), tier = A falls alle Auswahl-Instanzen der
Kategorie EIN Foto-Set teilen, sonst B. inclusion aus cnt.

Idempotent (TRUNCATE + reload).
Usage: python3 db_load_static.py [--src <pptx>]
"""
import argparse
import glob
import json
import os
import sys
from collections import defaultdict

import psycopg2
from psycopg2.extras import execute_values
from pptx import Presentation

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from _deckpipe import slugify                                  # noqa
from analyze_structure import deck_elements, headline           # noqa
from build_category_samples import full_text, img_sig           # noqa
from dedup_exact import CACHE                                   # noqa

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DATA = os.path.join(ROOT, "data")
CORPUS = "/Users/janrudat/Nextcloud/Kochfabrik Dokumente/AKARA_Präsentationen"
DSN = dict(host="localhost", port=5434, user="postgres",
           password="pptxgen", dbname="pptxgen")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--src", default=os.path.join(DATA,
                                                  "category_samples.pptx"))
    a = ap.parse_args()

    # 1) Jans Auswahl lesen (deck::page-Notizen)
    sel = []                                   # (slug,page,cat,ft,img)
    for s in Presentation(a.src).slides:
        n = (s.notes_slide.notes_text_frame.text.strip()
             if s.has_notes_slide else "")
        if "::" not in n:
            continue
        slug, _, pg = n.rpartition("::")
        pg = int(pg)
        cp = os.path.join(CACHE, slug, "elements.json")
        if not os.path.isfile(cp):
            continue
        seq = json.load(open(cp)).get(str(pg))
        if not seq:
            continue
        sel.append((slug, pg, headline(seq) or "?", full_text(seq),
                    img_sig(slug, seq)))
    if not sel:
        sys.exit(f"Keine Auswahl-Slides in {a.src}")
    wanted = {ft for _, _, _, ft, _ in sel}
    print(f"Auswahl: {len(sel)} Slides, "
          f"{len({c for _,_,c,_,_ in sel})} Kategorien")

    # 2) Korpus-Metadaten NUR für die Auswahl-Volltexte (cnt + Ø-pos)
    cx = psycopg2.connect(**DSN)
    cu = cx.cursor()
    cu.execute("SELECT deck, page FROM menu_composition")
    food = {(d, int(p)) for d, p in cu.fetchall()}
    cnt = defaultdict(int)
    pos = defaultdict(list)
    for pdf in sorted(glob.glob(os.path.join(CORPUS, "*.pdf"))):
        try:
            slug, el = deck_elements(pdf)
        except Exception:
            continue
        pages = sorted(int(k) for k in el if k != "_meta")
        n = len(pages)
        for idx, pg in enumerate(pages):
            if pg == 1 or (slug, pg) in food:
                continue
            ft = full_text(el[str(pg)])
            if ft in wanted:
                cnt[ft] += 1
                pos[ft].append(idx / max(n - 1, 1))

    # 3) je Kategorie gruppieren, golden = erste der Auswahl
    by_cat = defaultdict(list)
    for slug, pg, cat, ft, img in sel:
        by_cat[cat].append((slug, pg, ft, img))
    smap = {slugify(p): os.path.join(CORPUS, p)
            for p in os.listdir(CORPUS) if p.lower().endswith(".pdf")}
    cat_rank = sorted(by_cat,
                      key=lambda c: -max(cnt.get(ft, 0)
                                         for _, _, ft, _ in by_cat[c]))
    rows = []
    for rank, cat in enumerate(cat_rank, 1):
        inst = by_cat[cat]
        dom_ft = max((ft for _, _, ft, _ in inst), key=lambda f: cnt.get(f, 0))
        c = cnt.get(dom_ft, len(inst))
        sp = round(sum(pos[dom_ft]) / len(pos[dom_ft]), 4) if pos[dom_ft] \
            else None
        tier = "A" if len({im for _, _, _, im in inst}) <= 1 else "B"
        incl = ("pflicht" if c >= 100 else
                "bedingt" if c >= 10 else "optional")
        for i, (slug, pg, ft, img) in enumerate(inst):
            rows.append((cat[:80], rank, c, tier, sp, incl, slug,
                         smap.get(slug, ""), pg, ft, i == 0))

    cu.execute("TRUNCATE static_slide RESTART IDENTITY")
    execute_values(cu,
                   "INSERT INTO static_slide (category,rank,cnt,tier,"
                   "skel_pos,inclusion,deck,src_pdf,page,full_text,"
                   "is_golden) VALUES %s", rows)
    cx.commit()
    cu.execute("SELECT rank,category,cnt,tier,inclusion,"
               "round(skel_pos::numeric,2),"
               "count(*) FILTER (WHERE is_golden),"
               "count(*) FILTER (WHERE NOT is_golden) "
               "FROM static_slide GROUP BY 1,2,3,4,5,6 ORDER BY rank")
    print(f"\n{len(rows)} Zeilen aus deiner Auswahl (kein Korpus-Artefakt):")
    print("R  cnt T inclusion  pos  golden+alt  Kategorie")
    for rk, cat, c, t, inc, sp, g, al in cu.fetchall():
        print(f"{rk:>2} {c:>4} {t} {inc:9} {str(sp):>4}  "
              f"{g}g+{al}a   {cat[:38]}")
    cx.close()


if __name__ == "__main__":
    main()
