"""curate.py — regelbasiertes Slide-Löschen auf der Menü-Sammel-pptx.

Pro Slide wird die HEADLINE bestimmt (Text mit der größten Schrift) und
gegen eine Blockliste geprüft. Trifft eine Blockliste-Phrase die Headline
(Default) bzw. irgendeinen Text (--anywhere), fliegt die Slide raus.
Explizite Slide-Nummern via --drop. Notizen ("deck::page") bleiben
erhalten: gelöschte Junk-Slide nimmt ihre Notiz mit, behaltene behalten
ihre — die überlebende Notiz-Menge ist die menu_composition-Ground-Truth.

Idempotent + iterierbar: Phrase finden → --block ergänzen → erneut laufen.
Helper (norm/slide_texts) sind importierbar; CLI nur unter __main__.

Usage (aus spike-pptxgenjs/ oder beliebig):
  python3 ../scripts/curate.py /tmp/all_menus.pptx \
      --block "WERTSCHÄTZUNG IST DER SCHLÜSSEL" --block "PERSONAL"
  python3 ../scripts/curate.py /tmp/all_menus.pptx --drop 12,40,41 --dry
"""
import argparse
import os
import re
import sys

from pptx import Presentation


def norm(s):
    return re.sub(r"\s+", " ", (s or "")).strip().upper()


def slide_texts(slide):
    """(headline, fulltext): Headline = Text des Shapes mit größter Schrift."""
    full, best_sz, best_txt = [], -1.0, ""
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        txt = sh.text_frame.text
        if not txt.strip():
            continue
        full.append(txt)
        smax = 0.0
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size is not None:
                    smax = max(smax, float(r.font.size))
        if smax > best_sz:
            best_sz, best_txt = smax, txt
    return norm(best_txt), norm(" ".join(full))


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("inp", nargs="?", default="/tmp/all_menus.pptx")
    ap.add_argument("--out", default=None, help="Default: in-place (atomar)")
    ap.add_argument("--block", action="append", default=[],
                    help="Phrase; löscht Slide wenn in Headline (wiederholbar)")
    ap.add_argument("--drop", default="", help="explizite Slide-Nrn: 12,40,41")
    ap.add_argument("--anywhere", action="store_true",
                    help="Phrase irgendwo im Text statt nur Headline")
    ap.add_argument("--dry", action="store_true",
                    help="nur Report, nicht löschen")
    a = ap.parse_args()

    blocks = [norm(b) for b in a.block if b.strip()]
    drop_nos = {int(x) for x in re.findall(r"\d+", a.drop)}
    prs = Presentation(a.inp)
    slides = list(prs.slides)
    total = len(slides)

    hits = {b: 0 for b in blocks}
    to_del = []          # (1-based slide-nr, grund)
    for i, sl in enumerate(slides, 1):
        if i in drop_nos:
            to_del.append((i, "drop"))
            continue
        head, full = slide_texts(sl)
        hay = full if a.anywhere else head
        for b in blocks:
            if b and b in hay:
                hits[b] += 1
                to_del.append((i, f"block:{b[:40]}"))
                break

    print(f"Slides gesamt: {total}")
    print(f"Treffer: {len(to_del)}  →  bleiben: {total - len(to_del)}")
    for b in blocks:
        print(f"  block '{b[:50]}': {hits[b]}")
    if drop_nos:
        seen = sorted(drop_nos & {n for n, _ in to_del})
        print(f"  drop (explizit): {seen}")
    if to_del[:8]:
        ex = ", ".join(f"S{n}" for n, _ in to_del[:8])
        print(f"  Beispiele: {ex}{' …' if len(to_del) > 8 else ''}")

    if a.dry:
        print("\n[dry] nichts geschrieben.")
        return
    if not to_del:
        print("\nNichts zu löschen.")
        return

    # Löschen über sldIdLst (Notizen behaltener Slides bleiben unberührt).
    lst = prs.slides._sldIdLst
    ids = list(lst)
    for i, _ in sorted(to_del, reverse=True):
        lst.remove(ids[i - 1])

    out = a.out or a.inp
    tmp = out + ".tmp"
    prs.save(tmp)
    os.replace(tmp, out)

    chk = Presentation(out)
    withnote = sum(1 for s in chk.slides if s.has_notes_slide
                   and s.notes_slide.notes_text_frame.text.strip())
    print(f"\nOK → {out}")
    print(f"Slides jetzt: {len(list(chk.slides))} | "
          f"mit deck::page-Notiz: {withnote}")


if __name__ == "__main__":
    main()
