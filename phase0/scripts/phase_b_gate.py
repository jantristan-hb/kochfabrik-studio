"""phase_b_gate.py — Phase-B Mess-Gate.

Misst die Engine-Fehlerrate über eine STRATIFIZIERTE Stichprobe des
KOCHfabrik-Korpus. Liefert die Daten, um Phase C (Korpus-Härtung) zu
schätzen — KEINE Härtung hier.

Pro Deck: convert.py laufen lassen + Slide-1-Render Original vs. Recon
diffen + strukturelle Heuristiken (Text/Bild/Frame vorhanden). Fehler in
Klassen einsortieren. Aggregat → phase0/REPORT-phase-b.md.

Ausführen aus phase0/spike-pptxgenjs/:
  python3 ../scripts/phase_b_gate.py --n 25
"""
import argparse
import glob
import json
import os
import re
import subprocess
import sys
import tempfile

from PIL import Image
from pptx import Presentation
from pptx.util import Emu

CORPUS = "/Users/janrudat/Nextcloud/Kochfabrik Dokumente/AKARA_Präsentationen"
SPIKE = os.path.dirname(os.path.abspath(__file__)).replace(
    "/scripts", "/spike-pptxgenjs")
REPORT = os.path.abspath(os.path.join(SPIKE, "..", "REPORT-phase-b.md"))

STRATA = [
    ("foodkonzept|speisenidee|foodidee", "Foodkonzept"),
    ("eventkonzept|konzept", "Eventkonzept"),
    ("hochzeit|geburtstag|jubil", "Privatfeier"),
    ("bbq|sommerfest|weihnacht|festival", "Event/BBQ"),
    (".", "Kunden-benannt"),
]


def sample(n):
    pdfs = [p for p in sorted(glob.glob(os.path.join(CORPUS, "*.pdf")))
            if not os.path.basename(p).startswith("Angebot #")]
    per = max(1, n // len(STRATA))
    picked, used = [], set()
    for rx, _ in STRATA:
        hits = [p for p in pdfs if p not in used
                and re.search(rx, os.path.basename(p), re.I)]
        step = max(1, len(hits) // per) if hits else 1
        for p in hits[::step][:per]:
            picked.append(p); used.add(p)
    for p in pdfs:                       # auffüllen bis n
        if len(picked) >= n:
            break
        if p not in used:
            picked.append(p); used.add(p)
    return picked[:n]


def render_p1(pdf_or_pptx, out_png, work):
    src = pdf_or_pptx
    if src.endswith(".pptx"):
        subprocess.run(["soffice", "--headless", "--convert-to", "pdf",
                        "--outdir", work, src],
                       capture_output=True, timeout=120)
        src = os.path.join(work, os.path.splitext(
            os.path.basename(src))[0] + ".pdf")
    if not os.path.isfile(src):
        return False
    subprocess.run(["pdftoppm", "-png", "-r", "60", "-f", "1", "-l", "1",
                    src, out_png], capture_output=True, timeout=120)
    return os.path.isfile(out_png + "-1.png")


def gray_diff(a_png, b_png):
    try:
        a = Image.open(a_png).convert("L").resize((300, 170))
        b = Image.open(b_png).convert("L").resize((300, 170))
        pa, pb = list(a.getdata()), list(b.getdata())
        return sum(abs(x - y) for x, y in zip(pa, pb)) / (len(pa) * 255)
    except Exception:
        return 1.0


def classify(pptx):
    """Strukturelle Heuristiken → Fehlerklassen-Set."""
    flags = set()
    pr = Presentation(pptx)
    n = len(pr.slides)
    zero_text = 0
    for sl in pr.slides:
        runs = sum(len(p.runs) for sh in sl.shapes if sh.has_text_frame
                   for p in sh.text_frame.paragraphs)
        pics = sum(1 for sh in sl.shapes if sh.shape_type == 13)
        if runs == 0:
            zero_text += 1
        if pics == 0 and runs == 0:
            flags.add("leere-slide")
    if n and zero_text / n > 0.35:
        flags.add("text-fehlt")
    return flags, n


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--n", type=int, default=25)
    a = ap.parse_args()
    decks = sample(a.n)
    work = tempfile.mkdtemp(prefix="phaseb_")
    rows, classes = [], {}

    for i, pdf in enumerate(decks, 1):
        name = os.path.splitext(os.path.basename(pdf))[0]
        print(f"[{i}/{len(decks)}] {name[:48]}", file=sys.stderr)
        pptx = os.path.join(work, f"{i}.pptx")
        r = subprocess.run([sys.executable, os.path.join(SPIKE, "convert.py"),
                            pdf, pptx], cwd=SPIKE,
                           capture_output=True, text=True)
        if r.returncode != 0 or not os.path.isfile(pptx):
            msg = (r.stderr or r.stdout).strip().splitlines()[-1:][0:1]
            stg = re.search(r"\[(\w[\w-]*)\]", " ".join(msg))
            cls = f"pipeline:{stg.group(1) if stg else '?'}"
            rows.append((name, "FAIL", cls))
            classes[cls] = classes.get(cls, 0) + 1
            continue
        flags, nsl = classify(pptx)
        op, rp = os.path.join(work, f"o{i}"), os.path.join(work, f"r{i}")
        if render_p1(pdf, op, work) and render_p1(pptx, rp, work):
            d = gray_diff(op + "-1.png", rp + "-1.png")
            if d > 0.20:
                flags.add(f"visual-diff(>{0.20})")
        st = "OK" if not flags else "FLAG"
        rows.append((name, st, ",".join(sorted(flags)) or "-"))
        for f in flags:
            classes[f] = classes.get(f, 0) + 1

    total = len(rows)
    clean = sum(1 for _, s, _ in rows if s == "OK")
    fails = sum(1 for _, s, _ in rows if s == "FAIL")
    clean_rate = clean / total if total else 0
    if clean_rate >= 0.7:
        dec = ("**Phase C ≈ M** — Engine trägt breit; gezielte Fixes der "
               "Top-Fehlerklasse genügen.")
    elif clean_rate >= 0.4:
        dec = ("**Phase C ≈ L** — mehrere Fehlerklassen relevant; "
               "iterative Härtung über 2–3 Wellen.")
    else:
        dec = ("**Phase C ≈ XL / Strategie-Review** — Kalibrierung/Regeln "
               "generalisieren nicht; Ansatz pro Template-Familie prüfen.")

    md = [f"# REPORT — Phase-B Mess-Gate ({total} Decks, stratifiziert)\n",
          f"- Clean (OK): **{clean}/{total}** ({clean_rate*100:.0f}%)",
          f"- Flagged: {total-clean-fails} · Pipeline-FAIL: {fails}\n",
          "## Fehlerrate je Klasse\n",
          "| Klasse | Decks betroffen | Anteil |", "|---|--:|--:|"]
    for cls, c in sorted(classes.items(), key=lambda x: -x[1]):
        md.append(f"| {cls} | {c} | {c/total*100:.0f}% |")
    md += ["\n## Decks\n", "| Deck | Status | Klassen |", "|---|---|---|"]
    for nme, st, cl in rows:
        md.append(f"| {nme[:46]} | {st} | {cl} |")
    md += ["\n## Decision-Empfehlung\n", dec,
           "\n> Limitierung: Diff nur Slide 1, strukturelle Heuristik "
           "(keine Vision). Für Phase-C-Scoping ausreichend; exakte "
           "Fix-Liste entsteht beim Härten.\n"]
    open(REPORT, "w", encoding="utf-8").write("\n".join(md) + "\n")
    print(f"\nREPORT: {REPORT} | clean {clean}/{total} "
          f"({clean_rate*100:.0f}%) | fails {fails}")


if __name__ == "__main__":
    main()
