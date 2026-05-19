"""compose_offer.py — Angebot-PDF → passende Korpus-Kompositionen.

Phase D Kern (Eingabe-agnostisch, hier: KOCHfabrik-Angebot-PDF):
  match : Angebot in Gänge zerlegen → je Gang Top-K inhaltlich
          ähnlichste Slides aus dem kuratierten Korpus (Embedding
          headline+body = Gerichte). NUR Report — Verifikations-Gate.
  build : ausgewählte (deck,page) via _deckpipe faithful zu EINEM
          editierbaren Deck (korrektes KOCHfabrik-Logo) zusammenbauen.

Korpus = überlebende Notizen der kuratierten /tmp/all_menus.pptx
(Ground-Truth) ⨝ Texte aus /tmp/all_menus.slides.json.

Usage:
  python3 compose_offer.py match  "<angebot.pdf>"
  python3 compose_offer.py build  "<angebot.pdf>" --pick d::p,d::p,... [-o out.pptx]
"""
import argparse
import json
import os
import re
import subprocess
import sys
import tempfile

import numpy as np

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from _deckpipe import slugify, cached_deck                       # noqa

CORPUS_DIR = "/home/jrudat/Nextcloud/Kochfabrik Dokumente/AKARA_Präsentationen"
_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DATA = os.path.join(_ROOT, "data")                 # persistent, raus aus /tmp
CURATED = os.path.join(DATA, "all_menus.pptx")
SLIDES_JSON = os.path.join(DATA, "all_menus.slides.json")
SPIKE = os.path.join(_ROOT, "spike-pptxgenjs")
MODEL, DIM, TASK = "gemini-embedding-001", 768, "SEMANTIC_SIMILARITY"
DSN = dict(host="localhost", port=5434, user="postgres",
           password="pptxgen", dbname="pptxgen")

FOOTER = re.compile(r"(KOCHfabrik|Kochfabrik|koch-fabrik|Prisdorf|Pinneberg|"
                    r"BIC:|DE\d|HRB|Steuernummer|Gerichtsstand|www\.|@|"
                    r"Bankverbindung|Geschäftsführer|Gesellschafter)")
# Gang-Überschriften: GROSSBUCHSTABEN-Zeile, keine Preis-/Mengenzeile
HEAD = re.compile(r"^[A-ZÄÖÜ][A-ZÄÖÜ0-9 &/'\-\.]{3,40}$")
SKIP_HEAD = {"ANGEBOT", "PERSONAL", "TECHNIK", "LOGISTIK", "PROJEKT",
             "EVENT AUSSTATTUNG", "PORZELLAN", "PROJEKTMANAGEMENT"}


def _key():
    for ln in open(os.path.expanduser("~/work/.env")):
        if ln.startswith("GEMINI_API_KEY="):
            return ln.split("=", 1)[1].strip().strip('"')
    sys.exit("GEMINI_API_KEY fehlt")


def embed(texts):
    import time
    import urllib.request
    url = (f"https://generativelanguage.googleapis.com/v1beta/"
           f"models/{MODEL}:batchEmbedContents?key={_key()}")
    out = []
    for i in range(0, len(texts), 100):
        body = {"requests": [
            {"model": f"models/{MODEL}",
             "content": {"parts": [{"text": t or " "}]},
             "taskType": TASK, "outputDimensionality": DIM}
            for t in texts[i:i + 100]]}
        for att in range(4):
            try:
                req = urllib.request.Request(
                    url, json.dumps(body).encode(),
                    {"Content-Type": "application/json"})
                res = json.loads(urllib.request.urlopen(req, timeout=120)
                                 .read())
                break
            except Exception:
                if att == 3:
                    raise
                time.sleep(2 * (att + 1))
        out.extend(e["values"] for e in res["embeddings"])
        print(f"  embed {min(i+100,len(texts))}/{len(texts)}", file=sys.stderr)
    return np.asarray(out, dtype=np.float64)


MD_SKIP = {"VERANSTALTUNGSINFORMATIONEN", "VEREINBARUNG"}


def _parse_md(path, offer):
    """Markdown-Fixture: '## Angebot' nach <offer> wählen, '### GANG'-
    Header sammeln, Code-Fences/Tabellen/Briefkopf überspringen."""
    lines = open(path, encoding="utf-8").read().splitlines()
    # gewünschte Offer-Sektion ('## ' bis nächstes '## ' oder '---') slicen
    sec, take = [], False
    for ln in lines:
        if ln.startswith("## "):
            take = offer.lower() in ln.lower()
        elif ln.startswith("# ") or ln.strip() == "---":
            if take:
                break
            take = False
        if take:
            sec.append(ln)
    courses, cur, buf, fence = [], None, [], False
    for raw in sec:
        ln = raw.strip()
        if ln.startswith("```"):
            fence = not fence
            continue
        if fence or not ln or ln.startswith("|") or ln.startswith(">"):
            continue
        if ln.startswith("### "):
            title = ln[4:].strip()
            if cur and buf:
                courses.append((cur, " ".join(buf)))
            cur, buf = None, []
            if title.upper() not in MD_SKIP:
                cur = title
        elif cur:
            buf.append(ln)
    if cur and buf:
        courses.append((cur, " ".join(buf)))
    return [(c, b) for c, b in courses if len(b) > 25]


def parse_offer(path, offer=""):
    if path.lower().endswith((".md", ".txt")):
        return _parse_md(path, offer or "Angebot")
    txt = subprocess.run(["pdftotext", "-layout", path, "-"],
                          capture_output=True, text=True).stdout
    courses, cur, buf = [], None, []
    for raw in txt.splitlines():
        ln = raw.strip()
        if not ln or FOOTER.search(ln) or ln.isdigit():
            continue
        if HEAD.match(ln) and not re.search(r"\d", ln):
            key = ln.upper()
            if cur and buf:                       # laufenden Gang flushen
                courses.append((cur, " ".join(buf)))
            cur, buf = None, []
            if any(s in key for s in SKIP_HEAD):
                continue
            cur = ln
        elif cur:
            buf.append(ln)
    if cur and buf:
        courses.append((cur, " ".join(buf)))
    # nur Gänge mit echtem Speisen-Inhalt
    return [(c, b) for c, b in courses if len(b) > 25]


def load_corpus():
    from pptx import Presentation
    keep = set()
    for s in Presentation(CURATED).slides:
        if s.has_notes_slide:
            n = s.notes_slide.notes_text_frame.text.strip()
            if "::" in n:
                keep.add(n)
    rows = [r for r in json.load(open(SLIDES_JSON))
            if f"{r['deck']}::{r['page']}" in keep]
    smap = {slugify(p): os.path.join(CORPUS_DIR, p)
            for p in os.listdir(CORPUS_DIR) if p.lower().endswith(".pdf")}
    for r in rows:
        r["src"] = smap.get(r["deck"])
    return [r for r in rows if r["src"]]


def cos_topk(q, M, k):
    q = q / (np.linalg.norm(q) + 1e-9)
    Mn = M / (np.linalg.norm(M, axis=1, keepdims=True) + 1e-9)
    s = Mn @ q
    idx = np.argsort(-s)[:k]
    return [(int(i), float(s[i])) for i in idx]


def cmd_match(pdf, k, offer=""):
    """Produktiv: nur die Angebot-Gänge live embedden, Korpus-Treffer
    via pgvector-ANN (menu_composition.embedding, cosine). Kein
    Live-Re-Embed der 1010 mehr."""
    import psycopg2
    courses = parse_offer(pdf, offer)
    cx = psycopg2.connect(**DSN)
    cu = cx.cursor()
    cu.execute("SELECT count(*) FROM menu_composition "
               "WHERE embedding IS NOT NULL")
    n = cu.fetchone()[0]
    print(f"Angebot: {len(courses)} Speisen-Gänge | "
          f"Korpus (pgvector): {n} Slides")
    ce = embed([f"{c} — {b}" for c, b in courses])
    for (c, b), qv in zip(courses, ce):
        q = "[" + ",".join(f"{x:.6f}" for x in qv) + "]"
        cu.execute(
            "SELECT deck,page,headline,body,module_label,"
            "1-(embedding<=>%s::vector) sim FROM menu_composition "
            "ORDER BY embedding<=>%s::vector LIMIT %s", (q, q, k))
        print(f"\n### {c}  ({b[:70]}…)")
        for deck, pg, hl, bd, lab, sim in cu.fetchall():
            print(f"  {sim:.3f}  {deck[:26]}::{pg:>2}  "
                  f"[{(lab or '')[:16]:16}] {hl[:24]:24} | {bd[:42]}")
    cx.close()


def cmd_build(pdf, picks, out):
    rows = {f"{r['deck']}::{r['page']}": r for r in load_corpus()}
    sel = [rows[p] for p in picks if p in rows]
    if not sel:
        sys.exit("Keine gültigen Picks (Format deck::page).")
    shared = tempfile.mkdtemp(prefix="offer_")
    combined, notes, logos, meta = {}, {}, {}, None
    by_src = {}
    for r in sel:
        by_src.setdefault(r["src"], []).append(int(r["page"]))
    n = 0
    for src, pages in by_src.items():
        slug, el, lg = cached_deck(src, shared)
        logos.update(lg)
        if meta is None:
            meta = el.get("_meta", {"w_pt": 960, "h_pt": 540})
        for pg in pages:
            seq = el.get(str(pg))
            if not seq:
                print(f"  warn: {slug}::{pg} fehlt", file=sys.stderr)
                continue
            n += 1
            combined[str(n)] = seq
            notes[str(n)] = f"{slug}::{pg}"
    mm = dict(meta or {"w_pt": 960, "h_pt": 540})
    mm["deck"], mm["notes"] = "offer", notes
    combined["_meta"] = mm
    json.dump(combined, open(os.path.join(shared, "elements.json"), "w"))
    json.dump(logos, open(os.path.join(shared, "logos.json"), "w"))
    subprocess.run(["node", os.path.join(SPIKE, "reconstruct.js"),
                    "elements.json", out], cwd=shared,
                   capture_output=True, check=True, timeout=600)
    print(f"OK: {out} — {n} Slides aus {len(by_src)} Decks")


def parse_offer_dishes(path, offer=""):
    """Wie parse_offer, aber pro Gang die EINZEL-Gerichte (name, desc).
    md: Leerzeile trennt Gerichte; Zeile 1 = Name, Rest = Beschreibung."""
    if not path.lower().endswith((".md", ".txt")):
        # PDF-Fallback: ganzer Gang-Text als ein Block (grob, später feiner)
        return [(c, [(b, "")]) for c, b in parse_offer(path, offer)]
    lines = open(path, encoding="utf-8").read().splitlines()
    sec, take = [], False
    for ln in lines:
        if ln.startswith("## "):
            take = (offer or "Angebot").lower() in ln.lower()
        elif ln.startswith("# ") or ln.strip() == "---":
            if take:
                break
        if take:
            sec.append(ln)
    out, course, pend, fence = [], None, [], False
    skip = {"VERANSTALTUNGSINFORMATIONEN", "VEREINBARUNG"}

    def flush_dish():
        if course and pend:
            cdishes.append((pend[0], " ".join(pend[1:])[:140]))

    cdishes = []
    for raw in sec:
        ln = raw.strip()
        if ln.startswith("```"):
            fence = not fence
            continue
        if fence or ln.startswith("|") or ln.startswith(">"):
            continue
        if ln.startswith("### "):
            flush_dish()
            pend = []
            if course and cdishes:
                out.append((course, cdishes))
            cdishes = []
            t = ln[4:].strip()
            course = None if t.upper() in skip else t
        elif not ln:
            flush_dish()
            pend = []
        elif course:
            pend.append(ln)
    flush_dish()
    if course and cdishes:
        out.append((course, cdishes))
    return out


GOLD = "AA8339"


def text_swap(seq, dishes):
    """KF-Gericht-Captions sind separate Mini-Elemente: Bold-weiß = Name,
    folgende Regular = Beschreibung, Bold-gold = Sektions-Header (bleibt).
    In Lesereihenfolge (Spalten→y) zu (Name, [Desc])-Paaren gruppieren,
    NUR den Text der bestehenden Elemente ersetzen — Bbox + jeweiliges
    Eigen-Styling unangetastet → Position & Formatierung bleiben exakt."""
    texts = [e for e in seq if e["t"] == "text" and e.get("lines")]
    if not texts:
        return seq
    mx = max(max(l["size"] for l in e["lines"]) for e in texts)

    def w0(e):
        return str(e["lines"][0].get("weight", "")).lower()

    def c0(e):
        return e["lines"][0].get("color", "")

    def is_cap(e):                       # Menü-Größenbereich (nicht Headline)
        return max(l["size"] for l in e["lines"]) < 0.5 * mx
    def is_name(e):                      # Gericht-Name: Bold + weiß
        return w0(e) in ("bold", "extrabold") and c0(e) != GOLD
    def is_sect(e):                      # KF-Sektions-Header: Bold + gold
        return w0(e) in ("bold", "extrabold") and c0(e) == GOLD

    cap = sorted((e for e in texts if is_cap(e)),
                 key=lambda e: (round(e["x"] * 2) / 2, round(e["y"], 2)))
    pairs, i = [], 0                     # [(name_elem, [desc_elems])]
    while i < len(cap):
        if is_name(cap[i]):
            descs, j = [], i + 1
            while j < len(cap) and not is_name(cap[j]) \
                    and not is_sect(cap[j]):
                descs.append(cap[j])
                j += 1
            pairs.append((cap[i], descs))
            i = j
        else:
            i += 1

    def put(elem, txt):                  # Eigen-Stil halten, 1 Zeile, nur txt
        s = elem["lines"][0]
        elem["lines"] = [{k: s[k] for k in
                          ("size", "color", "weight", "italic")
                          if k in s} | {"txt": txt}]

    for k, (nm, descs) in enumerate(pairs):
        if k < len(dishes):
            name, desc = dishes[k]
            put(nm, name)
            if descs:
                put(descs[0], desc)
                for extra in descs[1:]:
                    put(extra, "")
        else:                            # überzählige Slots leeren
            put(nm, "")
            for d in descs:
                put(d, "")
    if len(dishes) > len(pairs) and pairs:   # Rest an letzte Desc hängen
        nm, descs = pairs[-1]
        tgt = descs[0] if descs else nm
        rest = " • ".join(f"{n} ({d})" if d else n
                          for n, d in dishes[len(pairs):])
        tgt["lines"][0]["txt"] = (tgt["lines"][0]["txt"] + " • "
                                  + rest).strip(" •")
    return seq


def cmd_swap(path, offer, picks, out):
    courses = parse_offer_dishes(path, offer)
    if not courses:
        sys.exit("Keine Gänge im Angebot geparst.")
    if len(picks) != len(courses):
        print(f"  warn: {len(picks)} Picks vs {len(courses)} Gänge — "
              f"zippe kürzeste", file=sys.stderr)
    pdfs = {slugify(p): os.path.join(CORPUS_DIR, p)
            for p in os.listdir(CORPUS_DIR) if p.lower().endswith(".pdf")}
    shared = tempfile.mkdtemp(prefix="swap_")
    combined, notes, logos, meta = {}, {}, {}, None
    n = 0
    for (course, dishes), pick in zip(courses, picks):
        slug, _, pg = pick.partition("::")
        src = pdfs.get(slug)
        if not src:
            print(f"  warn: PDF für {slug} fehlt", file=sys.stderr)
            continue
        _, el, lg = cached_deck(src, shared)
        logos.update(lg)
        if meta is None:
            meta = el.get("_meta", {"w_pt": 960, "h_pt": 540})
        seq = el.get(str(int(pg)))
        if not seq:
            print(f"  warn: {slug}::{pg} fehlt", file=sys.stderr)
            continue
        text_swap(seq, dishes)
        n += 1
        combined[str(n)] = seq
        notes[str(n)] = f"{course} ⇐ {slug}::{pg}"
        print(f"  {course:22} ⇐ {slug}::{pg}  ({len(dishes)} Gerichte)")
    mm = dict(meta or {"w_pt": 960, "h_pt": 540})
    mm["deck"], mm["notes"] = "offer-swap", notes
    combined["_meta"] = mm
    json.dump(combined, open(os.path.join(shared, "elements.json"), "w"))
    json.dump(logos, open(os.path.join(shared, "logos.json"), "w"))
    subprocess.run(["node", os.path.join(SPIKE, "reconstruct.js"),
                    "elements.json", out], cwd=shared,
                   capture_output=True, check=True, timeout=600)
    print(f"OK: {out} — {n} Slides, Text auf Angebot-Gerichte geswappt")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("cmd", choices=["match", "build", "swap"])
    ap.add_argument("pdf")
    ap.add_argument("-k", type=int, default=5)
    ap.add_argument("--pick", default="")
    ap.add_argument("--offer", default="",
                    help="Sektion in .md-Fixture (z.B. 'Nordlicht')")
    ap.add_argument("-o", default="/tmp/offer_deck.pptx")
    a = ap.parse_args()
    picks = [p.strip() for p in a.pick.split(",") if p.strip()]
    if a.cmd == "match":
        cmd_match(a.pdf, a.k, a.offer)
    elif a.cmd == "swap":
        cmd_swap(a.pdf, a.offer, picks, a.o)
    else:
        cmd_build(a.pdf, picks, a.o)


if __name__ == "__main__":
    main()
