"""dedup_exact.py — NUR wirklich eindeutige Duplikat-Slides entfernen.

Signatur je Slide (aus dem Cache via deck::page-Notiz):
  norm_text  = aller Text, whitespace/case-normalisiert
  geometry   = sortierte gerundete (typ,x,y,w,h) je Element
Zwei Slides = Duplikat NUR wenn norm_text UND geometry identisch.
Konservativ: nur Slides mit substantiellem Text (>=12 Zeichen) kommen
für Dedup in Frage — text-arme/Bild-only bleiben unangetastet
(event-spezifische unterscheiden sich im Text → bleiben automatisch).
Pro Duplikat-Gruppe bleibt die ERSTE (aktuelle Sortier-Position).

--dry (default): nur Report + Stichprobe. --apply: löschen (gehärtet,
Notiz-erhaltend, Waisen kompaktiert). Pristine-Backup existiert separat.

Usage: python3 dedup_exact.py /tmp/all_info.pptx [--apply]
"""
import argparse
import json
import os
import re
import sys
from collections import defaultdict

from pptx import Presentation
from pptx.oxml.ns import qn

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
CACHE = os.path.join(ROOT, "data", "cache")
SLIDE_RT = ("http://schemas.openxmlformats.org/officeDocument/2006/"
            "relationships/slide")


def sig(seq):
    txt, geo = [], []
    for e in seq:
        geo.append((e.get("t"), round(e.get("x", 0), 1),
                    round(e.get("y", 0), 1), round(e.get("w", 0), 1),
                    round(e.get("h", 0), 1)))
        if e.get("t") == "text":
            for ln in e.get("lines", []):
                txt.append(ln.get("txt", ""))
    nt = re.sub(r"\s+", " ", " ".join(txt)).strip().upper()
    return nt, (nt, tuple(sorted(geo)))


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx", nargs="?", default="/tmp/all_info.pptx")
    ap.add_argument("--apply", action="store_true")
    a = ap.parse_args()

    prs = Presentation(a.pptx)
    slides = list(prs.slides)
    groups = defaultdict(list)        # signature -> [slide_idx(1-based)]
    skipped = 0
    for i, s in enumerate(slides, 1):
        note = (s.notes_slide.notes_text_frame.text.strip()
                if s.has_notes_slide else "")
        if "::" not in note:
            skipped += 1
            continue
        slug, _, pg = note.rpartition("::")
        cp = os.path.join(CACHE, slug, "elements.json")
        if not os.path.isfile(cp):
            skipped += 1
            continue
        seq = json.load(open(cp)).get(str(int(pg)))
        if not seq:
            skipped += 1
            continue
        nt, signature = sig(seq)
        if len(nt) < 12:              # text-arm → NICHT dedupen
            skipped += 1
            continue
        groups[signature].append(i)

    dup_groups = {k: v for k, v in groups.items() if len(v) > 1}
    drop = sorted(i for v in dup_groups.values() for i in v[1:])
    print(f"Slides: {len(slides)} | dedup-fähig (Text>=12): "
          f"{sum(len(v) for v in groups.values())} | "
          f"text-arm/ohne Cache übersprungen: {skipped}")
    print(f"Eindeutige Duplikat-Gruppen: {len(dup_groups)} | "
          f"zu entfernen: {len(drop)} | bleiben: {len(slides)-len(drop)}")
    top = sorted(dup_groups.items(), key=lambda kv: -len(kv[1]))[:12]
    print("\nGrößte Duplikat-Gruppen (Kopien × — Text-Anfang):")
    for (nt, _), idxs in top:
        print(f"  {len(idxs):3d}×  {nt[:66]}")

    if not a.apply:
        print("\n[dry] nichts gelöscht. --apply zum Entfernen.")
        return
    if not drop:
        print("\nNichts zu entfernen.")
        return
    lst = prs.slides._sldIdLst
    ids = list(lst)
    for i in sorted(drop, reverse=True):
        sid = ids[i - 1]
        rid = sid.get(qn("r:id"))
        lst.remove(sid)
        try:
            prs.part.drop_rel(rid)
        except (KeyError, ValueError):
            pass
    live = {s.get(qn("r:id")) for s in lst}
    for rel in list(prs.part.rels.values()):
        if rel.reltype == SLIDE_RT and rel.rId not in live:
            prs.part.drop_rel(rel.rId)
    tmp = a.pptx + ".tmp"
    prs.save(tmp)
    os.replace(tmp, a.pptx)
    chk = Presentation(a.pptx)
    cs = list(chk.slides)
    wn = sum(1 for s in cs if s.has_notes_slide
             and s.notes_slide.notes_text_frame.text.strip())
    print(f"\nOK → {a.pptx} | {len(cs)} Slides, {wn} mit deck::page-Notiz")


if __name__ == "__main__":
    main()
